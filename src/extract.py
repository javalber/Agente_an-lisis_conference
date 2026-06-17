"""Extracción de texto plano — código determinista (NO usa modelo).

Todo se convierte a texto/markdown plano para que la capa de modelo (Fase 3)
sea agnóstica del formato de origen.

  - PDF:   pdfplumber (texto + tablas por página).
  - Excel: pandas + openpyxl (cada hoja -> tabla markdown).
  - PPTX:  python-pptx (texto de cada slide + NOTAS del presentador).
"""

from __future__ import annotations

import logging

from .state import FileKind

log = logging.getLogger(__name__)


def _table_to_markdown(table: list[list[str | None]]) -> str:
    """Convierte una tabla (lista de filas) extraída por pdfplumber a markdown."""
    rows = [[("" if cell is None else str(cell)).strip() for cell in row] for row in table]
    if not rows:
        return ""
    header = rows[0]
    body = rows[1:]
    md = ["| " + " | ".join(header) + " |"]
    md.append("| " + " | ".join("---" for _ in header) + " |")
    for r in body:
        # Normaliza el ancho de fila al del encabezado.
        r = (r + [""] * len(header))[: len(header)]
        md.append("| " + " | ".join(r) + " |")
    return "\n".join(md)


def _extract_pdf(path: str) -> str:
    import pdfplumber

    chunks: list[str] = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            chunks.append(f"== Página {i} ==")
            text = page.extract_text() or ""
            if text.strip():
                chunks.append(text)
            for ti, table in enumerate(page.extract_tables() or [], start=1):
                md = _table_to_markdown(table)
                if md:
                    chunks.append(f"-- Tabla {ti} --")
                    chunks.append(md)
    return "\n".join(chunks).strip()


def _extract_excel(path: str) -> str:
    import pandas as pd

    chunks: list[str] = []
    # sheet_name=None -> dict {nombre_hoja: DataFrame}
    sheets = pd.read_excel(path, sheet_name=None, engine="openpyxl")
    for name, df in sheets.items():
        chunks.append(f"## Hoja: {name}")
        if df.empty:
            chunks.append("(hoja vacía)")
        else:
            chunks.append(df.to_markdown(index=False))
    return "\n\n".join(chunks).strip()


def _extract_pptx(path: str) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    chunks: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        chunks.append(f"== Slide {i} ==")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if line:
                    chunks.append(line)
        # Notas del presentador: suelen traer el guion real.
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                chunks.append("-- Notas del presentador --")
                chunks.append(notes)
    return "\n".join(chunks).strip()


_EXTRACTORS = {
    "pdf": _extract_pdf,
    "excel": _extract_excel,
    "pptx": _extract_pptx,
}


def extract_text(path: str, kind: FileKind) -> str:
    """Despacha al extractor correspondiente. Lanza si el tipo no es soportado
    o si la extracción falla (para que el correo NO se marque como procesado)."""
    extractor = _EXTRACTORS.get(kind)
    if extractor is None:
        raise ValueError(f"Tipo de archivo no soportado para extracción: {kind!r}")
    text = extractor(path)
    log.info("Extraídos %d caracteres de %s (%s)", len(text), path, kind)
    return text
